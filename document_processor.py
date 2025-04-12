import os
import uuid
import random
from typing import List, Dict, Any
import pickle
import numpy as np
from pathlib import Path
import re
import requests
import json
import time

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_huggingface import HuggingFaceEmbeddings

# Document loaders
from pypdf import PdfReader
import docx2txt

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

try:
    from langdetect import detect, LangDetectException

    LANGDETECT_AVAILABLE = True
except ImportError:
    LANGDETECT_AVAILABLE = False
    print("Warning: langdetect not installed. Language detection will be limited.")

try:
    from googletrans import Translator

    GOOGLETRANS_AVAILABLE = True
except ImportError:
    GOOGLETRANS_AVAILABLE = False
    print("Warning: googletrans not installed. Translation will be limited.")


class DocumentProcessor:
    def __init__(self, persist_directory: str = "./doc_storage"):
        self.storage_dir = persist_directory
        self.documents = {}  # doc_id -> List[Document]
        self.document_embeddings = {}  # doc_id -> List[np.array]

        # Ollama configuration
        self.use_llm = True  # Set to False to disable LLM and use extractive summarization
        self.ollama_url = "http://localhost:11434/api/generate"
        self.llm_model = "llama3"  # Default model

        # Initialize Ollama for summarization
        self.available_models = self._get_available_ollama_models()
        if self.available_models:
            print(f"Available Ollama models: {', '.join(self.available_models)}")
            if self.llm_model not in self.available_models and self.available_models:
                self.llm_model = self.available_models[0]
                print(f"Using available model: {self.llm_model}")
        else:
            print("No Ollama models available. Using extractive summarization instead.")
            self.use_llm = False

        # Initialize translation components
        self.translator = None
        if GOOGLETRANS_AVAILABLE:
            try:
                self.translator = Translator()
                print("Google Translator initialized successfully")
            except Exception as e:
                print(f"Warning: Could not initialize Google Translator: {e}")

        # Lower the similarity threshold to show more results
        self.similarity_threshold = 0.2

        # Create embeddings model
        self.embeddings_model = None
        try:
            self.embeddings_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
            print("Embeddings model initialized successfully")
        except Exception as e:
            print(f"Warning: Could not initialize embeddings model: {e}")
            print("Search functionality will be limited.")

        # Create directory if not exists
        os.makedirs(self.storage_dir, exist_ok=True)
        os.makedirs(os.path.join(self.storage_dir, "translations"), exist_ok=True)
        os.makedirs(os.path.join(self.storage_dir, "summaries"), exist_ok=True)

        # Load existing documents if available
        self._load_documents()

    def _get_available_ollama_models(self) -> List[str]:
        """Get available Ollama models"""
        try:
            response = requests.get("http://localhost:11434/api/tags")
            if response.status_code == 200:
                data = response.json()
                if 'models' in data:
                    return [model["name"] for model in data["models"]]
            return []
        except Exception as e:
            print(f"Error getting Ollama models: {e}")
            return []

    def _load_documents(self):
        """Load documents from disk"""
        docs_file = os.path.join(self.storage_dir, "documents.pkl")
        if os.path.exists(docs_file):
            try:
                with open(docs_file, 'rb') as f:
                    self.documents = pickle.load(f)
                print(f"Loaded {len(self.documents)} documents from storage")
            except Exception as e:
                print(f"Error loading documents: {e}")
                self.documents = {}

    def _save_documents(self):
        """Save documents to disk"""
        try:
            docs_file = os.path.join(self.storage_dir, "documents.pkl")
            with open(docs_file, 'wb') as f:
                pickle.dump(self.documents, f)
            print(f"Saved {len(self.documents)} documents to storage")
        except Exception as e:
            print(f"Error saving documents: {e}")

    def detect_language(self, text: str) -> str:
        """Detect the language of text with improved Hindi detection"""
        try:
            # Handle empty text
            if not text or len(text.strip()) < 10:
                return "en"  # Default to English for very short text

            # First check for Devanagari script which is used for Hindi
            devanagari_chars = re.findall(r'[\u0900-\u097F\u0A00-\u0A7F\u0A80-\u0AFF]', text)

            # If there are a significant number of Devanagari characters, classify as Hindi
            if len(devanagari_chars) > len(text) * 0.1:  # More than 10% Devanagari characters
                print("Detected Hindi document based on Devanagari script")
                return "hi"

            # Use langdetect as a backup
            if LANGDETECT_AVAILABLE:
                try:
                    lang = detect(text)
                    if lang == 'hi':
                        return "hi"
                    return "en"  # Default to English for other languages
                except LangDetectException:
                    pass

            # Default to English if no Hindi detected
            return "en"

        except Exception as e:
            print(f"Language detection error: {e}")
            return "en"  # Default to English on error

    def translate_text(self, text: str, target_lang: str = "en", source_lang: str = None) -> str:
        """
        Translate text between languages with enhanced error handling
        """
        if not text or text.strip() == "":
            return text

        # Check if text is already in the target language
        if source_lang is None:
            try:
                source_lang = self.detect_language(text)
            except:
                source_lang = "en"  # Default

        # If already in target language, return as is
        if source_lang == target_lang:
            return text

        # Create translation cache key
        cache_key = f"{source_lang}_{target_lang}_{hash(text)}"
        cache_file = os.path.join(self.storage_dir, "translations", f"{cache_key}.txt")

        # Check if we have this translation cached
        if os.path.exists(cache_file):
            try:
                with open(cache_file, 'r', encoding='utf-8') as f:
                    return f.read()
            except:
                pass  # If cache read fails, continue with translation

        # Try Google Translate (most reliable)
        translated_text = None
        if GOOGLETRANS_AVAILABLE and self.translator:
            try:
                # For Hindi to English
                if source_lang == 'hi' and target_lang == 'en':
                    result = self.translator.translate(text, dest='en', src='hi')
                    if result and result.text:
                        translated_text = result.text
                # For English to Hindi
                elif source_lang == 'en' and target_lang == 'hi':
                    result = self.translator.translate(text, dest='hi', src='en')
                    if result and result.text:
                        translated_text = result.text
            except Exception as e:
                print(f"Google translation error: {e}")

        # If all translation methods fail, return original text
        if not translated_text:
            return text

        # Cache the translation for future use
        try:
            os.makedirs(os.path.join(self.storage_dir, "translations"), exist_ok=True)
            with open(cache_file, 'w', encoding='utf-8') as f:
                f.write(translated_text)
        except Exception as e:
            print(f"Error caching translation: {e}")

        return translated_text

    def process_document(self, file_path: str, file_name: str) -> str:
        """Process and store a document with enhanced language detection"""
        # Generate document ID
        doc_id = str(uuid.uuid4())

        # Parse document
        documents = self._parse_document(file_path, file_name, doc_id)

        # Detect document language
        doc_text = " ".join([doc.page_content for doc in documents])
        doc_lang = self.detect_language(doc_text)
        print(f"Detected document language: {doc_lang}")

        # Add language to document metadata
        for doc in documents:
            doc.metadata['language'] = doc_lang

            # For Hindi documents, add English translation to help with embeddings
            if doc_lang == 'hi':
                try:
                    # Translate chunks for better search
                    translated_text = self.translate_text(doc.page_content, 'en', 'hi')
                    doc.metadata['translated_content'] = translated_text
                except Exception as e:
                    print(f"Error translating document: {e}")

        # Split into chunks
        doc_chunks = self._split_documents(documents)

        # Store document chunks
        self.documents[doc_id] = doc_chunks

        # Compute embeddings with language awareness
        if self.embeddings_model:
            try:
                embeddings = []
                for chunk in doc_chunks:
                    chunk_text = chunk.page_content

                    # If document is in Hindi, use translated content for embedding
                    if doc_lang == 'hi' and 'translated_content' in chunk.metadata:
                        chunk_text = chunk.metadata['translated_content']

                    # Compute embedding
                    embedding = self.embeddings_model.embed_query(chunk_text)
                    embeddings.append(embedding)

                self.document_embeddings[doc_id] = embeddings
            except Exception as e:
                print(f"Warning: Could not compute embeddings: {e}")

        # Save to disk
        self._save_documents()

        return doc_id

    def _parse_document(self, file_path: str, file_name: str, doc_id: str) -> List[Document]:
        """Extract text from various document formats with improved error handling"""
        file_extension = os.path.splitext(file_name)[1].lower()

        metadata = {
            "source": file_name,
            "doc_id": doc_id
        }

        text_content = ""

        try:
            if file_extension == '.pdf':
                # Parse PDF with better handling
                reader = PdfReader(file_path)
                text_content = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        # Clean up strange characters and patterns
                        page_text = self._clean_text(page_text)
                        text_content += page_text + "\n\n"

            elif file_extension in ['.docx', '.doc']:
                # Parse Word document
                text_content = docx2txt.process(file_path)
                text_content = self._clean_text(text_content)

            elif file_extension in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
                # Parse PowerPoint
                try:
                    prs = Presentation(file_path)
                    text_content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content += shape.text + "\n\n"
                    text_content = self._clean_text(text_content)
                except Exception as e:
                    print(f"Error parsing PowerPoint: {e}")
                    text_content = f"Error parsing PowerPoint document: {str(e)}"

            else:
                # Try to read as text file
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        text_content = f.read()
                except UnicodeDecodeError:
                    # Try with different encoding if UTF-8 fails
                    with open(file_path, 'r', encoding='latin-1') as f:
                        text_content = f.read()
                except Exception as e:
                    print(f"Error reading file as text: {e}")
                    text_content = f"Error reading file: {str(e)}"

                text_content = self._clean_text(text_content)

        except Exception as e:
            print(f"Error parsing document {file_name}: {e}")
            text_content = f"Error parsing document: {str(e)}"

        # Create a document
        if text_content.strip():
            # Final clean up
            text_content = self._clean_text(text_content)
            return [Document(page_content=text_content, metadata=metadata)]
        else:
            # Return a placeholder if we couldn't extract text
            return [Document(page_content="No text could be extracted from this document.",
                             metadata=metadata)]

    def _clean_text(self, text: str) -> str:
        """Clean extracted text to remove strange patterns and characters, preserving Hindi"""
        if not text:
            return ""

        # Check if text contains Devanagari script (Hindi)
        if re.search(r'[\u0900-\u097F]', text):
            # For Hindi text, use more gentle cleaning to preserve characters
            # Remove non-printable characters and control characters
            text = re.sub(r'[\x00-\x1F\x7F-\x9F]+', ' ', text)

            # Replace multiple spaces, newlines with single ones
            text = re.sub(r'\s+', ' ', text)

            # Trim whitespace
            return text.strip()

        # For non-Hindi text, use the regular cleaning
        # Remove repeating patterns like "Date Page Date Page"
        text = re.sub(r'(Date\s+Page\s*)+', ' ', text, flags=re.IGNORECASE)
        text = re.sub(r'(Date\s+Pa\s*e\s*)+', ' ', text, flags=re.IGNORECASE)
        text = re.sub(r'(dæssmate|dæssptß|classmate|classp)\s*', ' ', text, flags=re.IGNORECASE)

        # Remove non-printable characters and control characters
        text = re.sub(r'[\x00-\x1F\x7F-\x9F]+', ' ', text)

        # Remove isolated special characters
        text = re.sub(r'(?<!\w)[+\-•#*=_]+(?!\w)', ' ', text)

        # Replace multiple spaces, newlines with single ones
        text = re.sub(r'\s+', ' ', text)

        # Trim whitespace
        text = text.strip()

        return text

    def _split_documents(self, documents: List[Document]) -> List[Document]:
        """Split documents into chunks for better processing"""
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )

        chunks = []
        for doc in documents:
            doc_chunks = text_splitter.split_documents([doc])
            chunks.extend(doc_chunks)

        return chunks

    def get_document(self, doc_id: str) -> List[Document]:
        """Get all chunks for a specific document"""
        return self.documents.get(doc_id, [])

    def get_document_images(self, doc_id: str) -> List[Dict[str, Any]]:
        """Create reliable document visualizations"""
        chunks = self.get_document(doc_id)
        if not chunks:
            return []

        # Generate images
        images = []
        doc_lang = chunks[0].metadata.get('language', 'en')
        doc_title = chunks[0].metadata.get('source', 'Document')

        # Extract key topics for visualization
        doc_content = "\n\n".join([chunk.page_content for chunk in chunks[:3]])
        topics = self._extract_visualization_topics(doc_content, doc_lang)

        # Use a reliable public image API for visualization
        for i in range(min(5, len(topics))):
            topic = topics[i]

            # Create a visual based on the topic
            image_info = self._create_topic_visual(doc_id, topic, i, doc_lang)
            if image_info:
                images.append(image_info)

        return images

    def _extract_visualization_topics(self, text: str, language: str) -> List[str]:
        """Extract document topics for visualization"""
        # For Hindi content, translate to English for processing
        if language == 'hi':
            try:
                text = self.translate_text(text, 'en', 'hi')
            except:
                pass

        # Basic topic extraction
        topics = []

        # Try to identify document type and extract relevant topics
        if "financial" in text.lower() or "revenue" in text.lower() or "profit" in text.lower():
            topics.append("Financial Analysis")

        if "research" in text.lower() or "study" in text.lower() or "findings" in text.lower():
            topics.append("Research Findings")

        if "project" in text.lower() or "timeline" in text.lower() or "milestone" in text.lower():
            topics.append("Project Timeline")

        if "strategy" in text.lower() or "plan" in text.lower() or "objective" in text.lower():
            topics.append("Strategic Plan")

        if "data" in text.lower() or "statistics" in text.lower() or "analysis" in text.lower():
            topics.append("Data Visualization")

        # Add some generic topics if we couldn't extract enough specific ones
        generic_topics = [
            "Document Overview",
            "Key Concepts",
            "Main Points",
            "Summary View",
            "Content Analysis"
        ]

        # Fill with generic topics if needed
        while len(topics) < 5:
            for topic in generic_topics:
                if topic not in topics:
                    topics.append(topic)
                    if len(topics) >= 5:
                        break

        # Translate topics to Hindi if needed
        if language == 'hi':
            hindi_topics = []
            hindi_translations = {
                "Financial Analysis": "वित्तीय विश्लेषण",
                "Research Findings": "शोध निष्कर्ष",
                "Project Timeline": "परियोजना समयरेखा",
                "Strategic Plan": "रणनीतिक योजना",
                "Data Visualization": "डेटा विज़ुअलाइज़ेशन",
                "Document Overview": "दस्तावेज़ अवलोकन",
                "Key Concepts": "मुख्य अवधारणाएँ",
                "Main Points": "मुख्य बिंदु",
                "Summary View": "सारांश दृश्य",
                "Content Analysis": "सामग्री विश्लेषण"
            }

            for topic in topics:
                if topic in hindi_translations:
                    hindi_topics.append(hindi_translations[topic])
                else:
                    try:
                        hindi_topics.append(self.translate_text(topic, 'hi', 'en'))
                    except:
                        hindi_topics.append(topic)

            return hindi_topics

        return topics

    def _create_topic_visual(self, doc_id: str, topic: str, index: int, language: str) -> Dict[str, Any]:
        """Create a reliable visual based on a topic using public image APIs"""
        # Define color themes by index
        colors = [
            "4361ee",  # Blue
            "3a0ca3",  # Indigo
            "f72585",  # Pink
            "4cc9f0",  # Light Blue
            "7209b7"  # Purple
        ]
        color = colors[index % len(colors)]

        # Create a formatted placeholder URL with the topic
        safe_topic = topic.replace(' ', '+')

        # Use a more reliable public image service
        placeholder_url = f"https://placehold.co/800x450/{color}/FFFFFF?text={safe_topic}"

        # Create image metadata
        image_info = {
            'filename': f"{doc_id}_visual_{index + 1}.png",
            'page_num': index + 1,
            'width': 800,
            'height': 450,
            'nearby_text': topic,
            'doc_id': doc_id,
            'is_placeholder': True,
            'placeholder_url': placeholder_url,
            'url': placeholder_url,
            'type': "visualization",
            'topic': topic,
            'language': language
        }

        return image_info

    def search_documents(self, query: str, k: int = 5, query_language: str = None) -> List[Dict[str, Any]]:
        """Search for documents with multilingual support"""
        results = []

        # Detect query language if not provided
        if query_language is None:
            query_language = self.detect_language(query)

        print(f"Query language detected: {query_language}")

        # If we have no embeddings model, just return the most recent documents
        if not self.embeddings_model or not self.document_embeddings:
            # Return most recent documents
            recent_docs = list(self.documents.items())[-k:]
            for doc_id, chunks in recent_docs:
                if chunks:
                    doc = chunks[0]  # Use first chunk for metadata

                    # Handle translation for result snippet if needed
                    doc_lang = doc.metadata.get('language', 'en')
                    snippet = chunks[0].page_content[:200] + "..." if len(chunks[0].page_content) > 200 else chunks[
                        0].page_content

                    # Translate snippet if languages don't match
                    if doc_lang != query_language:
                        try:
                            snippet = self.translate_text(snippet, target_lang=query_language, source_lang=doc_lang)
                        except Exception as e:
                            print(f"Error translating snippet: {e}")

                    results.append({
                        'doc_id': doc_id,
                        'title': doc.metadata.get('source', 'Unknown'),
                        'content': chunks[0].page_content,
                        'snippet': snippet,
                        'score': 1.0,  # Placeholder score
                        'language': doc_lang
                    })
            return results

        # Prepare query for embedding - translate if needed
        query_for_embedding = query
        if query_language == 'hi':
            try:
                # Translate Hindi query to English for better embedding
                query_for_embedding = self.translate_text(query, 'en', 'hi')
                print(f"Translated query: {query_for_embedding}")
            except Exception as e:
                print(f"Error translating query for embedding: {e}")

        # Get query embedding
        try:
            query_embedding = self.embeddings_model.embed_query(query_for_embedding)
            query_embedding = np.array(query_embedding)

            # Calculate similarity for each document chunk
            doc_scores = []
            seen_docs = set()  # Track which documents we've already matched

            for doc_id, embeddings in self.document_embeddings.items():
                if not self.documents.get(doc_id):
                    continue

                chunks = self.documents[doc_id]
                doc_lang = chunks[0].metadata.get('language', 'en') if chunks else 'en'

                # Find best matching chunk
                best_score = -1
                best_chunk_idx = 0

                for i, embedding in enumerate(embeddings):
                    if i >= len(chunks):
                        continue

                    embedding = np.array(embedding)

                    # Compute cosine similarity
                    similarity = np.dot(query_embedding, embedding) / (
                            np.linalg.norm(query_embedding) * np.linalg.norm(embedding)
                    )

                    if similarity > best_score:
                        best_score = similarity
                        best_chunk_idx = i

                # Only include if it passes the threshold and we haven't seen this doc yet
                if best_score > self.similarity_threshold and best_chunk_idx < len(chunks) and doc_id not in seen_docs:
                    doc_scores.append((doc_id, chunks[best_chunk_idx], best_score, doc_lang))
                    seen_docs.add(doc_id)

            # Sort by score
            doc_scores.sort(key=lambda x: x[2], reverse=True)

            # Format results
            for doc_id, chunk, score, doc_lang in doc_scores[:k]:
                # Get content and snippet
                content = chunk.page_content
                snippet = content[:200] + "..." if len(content) > 200 else content

                # Translate if languages don't match
                if doc_lang != query_language:
                    try:
                        snippet = self.translate_text(snippet, target_lang=query_language, source_lang=doc_lang)
                    except Exception as e:
                        print(f"Error translating result: {e}")

                results.append({
                    'doc_id': doc_id,
                    'title': chunk.metadata.get('source', 'Unknown'),
                    'content': content,
                    'snippet': snippet,
                    'score': float(score),
                    'language': doc_lang
                })

        except Exception as e:
            print(f"Error during search: {e}")
            # Fallback to returning recent documents
            recent_docs = list(self.documents.items())[-k:]
            for doc_id, chunks in recent_docs:
                if chunks:
                    doc = chunks[0]  # Use first chunk for metadata
                    doc_lang = doc.metadata.get('language', 'en')
                    snippet = chunks[0].page_content[:200] + "..." if len(chunks[0].page_content) > 200 else chunks[
                        0].page_content

                    # Translate snippet if languages don't match
                    if doc_lang != query_language:
                        try:
                            snippet = self.translate_text(snippet, target_lang=query_language, source_lang=doc_lang)
                        except Exception as e:
                            print(f"Error translating fallback snippet: {e}")

                    results.append({
                        'doc_id': doc_id,
                        'title': doc.metadata.get('source', 'Unknown'),
                        'content': chunks[0].page_content,
                        'snippet': snippet,
                        'score': 1.0,  # Placeholder score
                        'language': doc_lang
                    })

        return results

    def get_similar_documents(self, query: str, current_doc_id: str = None, k: int = 3, query_language: str = None) -> \
    List[Dict[str, Any]]:
        """Get documents similar to the query, excluding the current document"""
        results = self.search_documents(query, k=k + 3, query_language=query_language)

        # Filter out current document
        if current_doc_id:
            results = [r for r in results if r['doc_id'] != current_doc_id]

        return results[:k]

    def _generate_llm_summary(self, content: str, language: str = "en") -> str:
        """Generate a summary using LLM"""
        if not self.use_llm or not self.available_models:
            return self._generate_text_summary(content)

        # Check if content is too short for summarization
        if len(content.split()) < 50:
            return content

        # Create cache key for the summary
        cache_key = f"summary_{hash(content)}_{language}"
        cache_file = os.path.join(self.storage_dir, "summaries", f"{cache_key}.txt")

        # Check if we have this summary cached
        if os.path.exists(cache_file):
            try:
                with open(cache_file, 'r', encoding='utf-8') as f:
                    return f.read()
            except:
                pass  # If cache read fails, continue with generation

        # Limit content length to avoid prompt limits
        max_tokens = 4000
        content_words = content.split()
        if len(content_words) > max_tokens:
            content = " ".join(content_words[:max_tokens])
            content += "..."

        # Create prompt based on language
        if language == "hi":
            prompt = f"""
निम्नलिखित दस्तावेज़ का एक व्यापक, संरचित और सटीक सारांश लिखें:

{content}

एक पूर्ण, विस्तृत सारांश प्रदान करें जो दस्तावेज़ की मुख्य बातों को कवर करता है। आरंभिक परिचय, मुख्य विचारों का विवरण और निष्कर्ष/प्रमुख बिंदुओं को जोड़ना सुनिश्चित करें। कृपया हिंदी में उत्तर दें।

सारांश:
"""
        else:
            prompt = f"""
Write a comprehensive, structured, and accurate summary of the following document:

{content}

Provide a complete, detailed summary that covers the main points of the document. Be sure to include an initial introduction, details of the main ideas, and include the conclusion/key takeaways.

Summary:
"""

        try:
            # Call Ollama API for inference
            response = requests.post(
                self.ollama_url,
                json={
                    "model": self.llm_model,
                    "prompt": prompt,
                    "stream": False,
                    "options": {
                        "temperature": 0.1,  # Lower temperature for more factual output
                        "num_predict": 2048,  # Allow for longer responses
                        "stop": ["</response>", "</summary>"]  # Stop tokens
                    }
                },
                timeout=60  # Allow up to 60 seconds for response
            )

            if response.status_code == 200:
                # Extract the generated text
                result = response.json()
                summary = result.get("response", "")

                # Cleanup the summary - remove any starting/ending quotes or unnecessary prefixes
                summary = re.sub(r'^[\s"\'#*]+', '', summary)
                summary = re.sub(r'[\s"\']+$', '', summary)

                # Remove any "Summary:" prefix that might be in the output
                summary = re.sub(r'^(Summary|सारांश):?\s*', '', summary, flags=re.IGNORECASE)

                # Cache the summary
                try:
                    os.makedirs(os.path.join(self.storage_dir, "summaries"), exist_ok=True)
                    with open(cache_file, 'w', encoding='utf-8') as f:
                        f.write(summary)
                except Exception as e:
                    print(f"Error caching summary: {e}")

                return summary
            else:
                print(f"Error from Ollama API: {response.status_code} - {response.text}")
                return self._generate_text_summary(content)

        except Exception as e:
            print(f"Error generating LLM summary: {e}")
            return self._generate_text_summary(content)

    def _extract_llm_key_points(self, content: str, query: str = "", language: str = "en", count: int = 5) -> List[str]:
        """Extract key points using LLM"""
        if not self.use_llm or not self.available_models:
            return self._extract_key_points(content, count)

        # Check if content is too short for key point extraction
        if len(content.split()) < 50:
            sentences = content.split('.')
            return [s.strip() + '.' for s in sentences[:count] if s.strip()]

        # Create cache key for the key points
        cache_query = query[:100] if query else "no_query"
        cache_key = f"keypoints_{hash(content)}_{hash(cache_query)}_{language}_{count}"
        cache_file = os.path.join(self.storage_dir, "summaries", f"{cache_key}.json")

        # Check if we have these key points cached
        if os.path.exists(cache_file):
            try:
                with open(cache_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except:
                pass  # If cache read fails, continue with generation

        # Limit content length to avoid prompt limits
        max_tokens = 4000
        content_words = content.split()
        if len(content_words) > max_tokens:
            content = " ".join(content_words[:max_tokens])
            content += "..."

        # Create prompt based on language and query
        if language == "hi":
            if query:
                prompt = f"""
निम्नलिखित दस्तावेज़ के आधार पर, इस प्रश्न के संबंध में ठीक {count} महत्वपूर्ण बिंदु निकालें: "{query}"

दस्तावेज़:
{content}

प्रत्येक बिंदु एक पूर्ण वाक्य होना चाहिए जो प्रश्न से संबंधित महत्वपूर्ण जानकारी देता है। अपने उत्तर को बुलेट पॉइंट (•) के रूप में सूचीबद्ध करें। कृपया हिंदी में उत्तर दें।

महत्वपूर्ण बिंदु:
"""
            else:
                prompt = f"""
निम्नलिखित दस्तावेज़ से ठीक {count} सबसे महत्वपूर्ण बिंदु निकालें:

{content}

प्रत्येक बिंदु एक पूर्ण वाक्य होना चाहिए जो दस्तावेज़ की एक मुख्य जानकारी या अवधारणा को संक्षेप में प्रस्तुत करता है। अपने उत्तर को बुलेट पॉइंट (•) के रूप में सूचीबद्ध करें। कृपया हिंदी में उत्तर दें।

महत्वपूर्ण बिंदु:
"""
        else:
            if query:
                prompt = f"""
Based on the following document, extract exactly {count} key points that are relevant to this query: "{query}"

Document:
{content}

Each key point should be a complete sentence that provides important information related to the query. 
List your response as bullet points. Number each point from 1 to {count}.

KEY POINTS:
"""
            else:
                prompt = f"""
Extract exactly {count} most important key points from the following document:

{content}

Each key point should be a complete sentence that concisely represents a main piece of information or concept from the document.
List your response as bullet points. Number each point from 1 to {count}.

KEY POINTS:
"""

        try:
            # Call Ollama API for inference
            response = requests.post(
                self.ollama_url,
                json={
                    "model": self.llm_model,
                    "prompt": prompt,
                    "stream": False,
                    "options": {
                        "temperature": 0.1,  # Lower temperature for more factual output
                        "num_predict": 1024,  # Allow for longer responses
                        "stop": ["</response>", "</key_points>"]  # Stop tokens
                    }
                },
                timeout=30  # Allow up to 30 seconds for response
            )

            if response.status_code == 200:
                # Extract the generated text
                result = response.json()
                response_text = result.get("response", "")

                # Process the response to extract individual points
                key_points = []

                # Try to extract numbered or bulleted points
                if "1." in response_text or "•" in response_text or "-" in response_text:
                    # Extract points with various formats (numbers, bullets, etc.)
                    point_matches = re.findall(r'(?:^|\n)(?:\d+\.|\*|•|-)\s*(.+?)(?=\n(?:\d+\.|\*|•|-|$)|$)',
                                               response_text)
                    key_points = [p.strip() for p in point_matches if p.strip()]

                # If no structured points found, split by newlines and try to clean up
                if not key_points:
                    raw_points = [p.strip() for p in response_text.split('\n') if p.strip()]

                    # Remove any "Key Points:" headers or similar
                    key_points = []
                    for point in raw_points:
                        # Skip headers or very short segments
                        if len(point) < 10 or re.match(r'^(Key Points|महत्वपूर्ण बिंदु|Points|बिंदु):?$', point,
                                                       re.IGNORECASE):
                            continue

                        # Remove numbering or bullets at the start
                        point = re.sub(r'^[\d\.\*•\-\s]+', '', point)

                        # Add if meaningful content
                        if len(point) > 10:
                            key_points.append(point)

                # Ensure we don't exceed the requested count
                key_points = key_points[:count]

                # Ensure each point ends with proper punctuation
                for i in range(len(key_points)):
                    if not key_points[i].endswith(('.', '!', '?')):
                        key_points[i] += '.'

                # Cache the key points
                try:
                    os.makedirs(os.path.join(self.storage_dir, "summaries"), exist_ok=True)
                    with open(cache_file, 'w', encoding='utf-8') as f:
                        json.dump(key_points, f)
                except Exception as e:
                    print(f"Error caching key points: {e}")

                return key_points
            else:
                print(f"Error from Ollama API: {response.status_code} - {response.text}")
                return self._extract_key_points(content, count)

        except Exception as e:
            print(f"Error extracting LLM key points: {e}")
            return self._extract_key_points(content, count)

    def _generate_text_summary(self, text: str) -> str:
        """Generate a text summary using extractive summarization"""
        try:
            # Split into sentences
            sentences = re.split(r'(?<=[.!?])\s+', text)

            if len(sentences) <= 5:
                return text

            # Score sentences based on position
            scored_sentences = []
            for i, sentence in enumerate(sentences):
                # Prioritize first and last sentences
                position_score = 0
                if i < len(sentences) * 0.2:  # First 20%
                    position_score = 0.8 - (i * 0.2 / len(sentences))
                elif i > len(sentences) * 0.8:  # Last 20%
                    position_score = 0.4 * ((i - len(sentences) * 0.8) / (len(sentences) * 0.2))
                else:
                    position_score = 0.2

                # Score based on sentence length (prefer medium length)
                length = len(sentence.split())
                length_score = 0
                if 5 <= length <= 25:
                    length_score = 0.5
                elif length < 5:
                    length_score = 0.1
                else:
                    length_score = 0.3

                # Final score
                score = position_score + length_score
                scored_sentences.append((sentence, score, i))

            # Sort by score and then by original position
            scored_sentences.sort(key=lambda x: (-x[1], x[2]))

            # Take top sentences (about 20% of original text)
            num_summary_sentences = max(5, int(len(sentences) * 0.2))
            selected_sentences = scored_sentences[:num_summary_sentences]

            # Sort back to original order
            selected_sentences.sort(key=lambda x: x[2])

            # Join into summary
            summary = " ".join([s[0] for s in selected_sentences])

            return summary

        except Exception as e:
            print(f"Error generating summary: {e}")
            # Fallback to first 3 sentences
            sentences = text.split('.')[:3]
            return '. '.join(sentences) + '.'

    def _extract_key_points(self, text: str, count: int = 5) -> List[str]:
        """Extract key points from text"""
        try:
            # Split into sentences
            sentences = re.split(r'(?<=[.!?])\s+', text)

            if len(sentences) <= count:
                return sentences

            # Score sentences based on position and content
            scored_sentences = []

            # Get common terms
            words = re.findall(r'\b\w+\b', text.lower())
            word_freq = {}

            # Common stopwords to ignore
            stop_words = {'a', 'an', 'the', 'and', 'or', 'but', 'is', 'are', 'was', 'were',
                          'to', 'of', 'in', 'for', 'with', 'by', 'at', 'on'}

            for word in words:
                if word not in stop_words and len(word) > 3:
                    word_freq[word] = word_freq.get(word, 0) + 1

            # Score sentences based on term frequency
            for i, sentence in enumerate(sentences):
                score = 0
                sentence_words = re.findall(r'\b\w+\b', sentence.lower())

                # Score based on important terms
                for word in sentence_words:
                    if word in word_freq:
                        score += word_freq[word]

                # Normalize by sentence length to avoid favoring long sentences
                if len(sentence_words) > 0:
                    score = score / len(sentence_words) ** 0.5

                # Boost for sentences with numbers or dates (often important facts)
                if re.search(r'\d', sentence):
                    score *= 1.2

                scored_sentences.append((sentence, score))

            # Sort by score
            scored_sentences.sort(key=lambda x: x[1], reverse=True)

            # Take top sentences
            key_points = [s for s, _ in scored_sentences[:count]]

            return key_points

        except Exception as e:
            print(f"Error extracting key points: {e}")
            # Fallback to first few sentences
            sentences = text.split('.')[:count]
            return [s.strip() + '.' for s in sentences if s.strip()]

    def get_document_summary(self, doc_id: str, target_language: str = None, use_ai: bool = True) -> Dict[str, Any]:
        """
        Get a summary of the document in the requested language with better Hindi support
        Uses LLM for summarization when available
        """
        # Get document chunks
        chunks = self.get_document(doc_id)
        if not chunks:
            return {
                "summary": "Document not found",
                "key_points": [],
                "language": target_language or "en"
            }

        # Get document language
        doc_lang = chunks[0].metadata.get('language', 'en')

        # If target language not specified, use document language
        if target_language is None:
            target_language = doc_lang

        # Get document content (first 5 chunks or all if fewer)
        content = "\n\n".join([chunk.page_content for chunk in chunks[:5]])

        # Check if LLM summarization is requested and available
        if use_ai and self.use_llm and self.available_models:
            # For Hindi documents, we need to ensure the summary is also in Hindi
            if doc_lang == 'hi':
                print("Generating Hindi summary for Hindi document")

                # If target language is Hindi, generate directly in Hindi
                if target_language == 'hi':
                    summary = self._generate_llm_summary(content, language='hi')
                    key_points = self._extract_llm_key_points(content, language='hi')
                else:
                    # Generate in Hindi then translate to target language
                    summary = self._generate_llm_summary(content, language='hi')
                    key_points = self._extract_llm_key_points(content, language='hi')

                    # Translate to target language
                    try:
                        summary = self.translate_text(summary, target_lang=target_language, source_lang='hi')
                        translated_points = []
                        for point in key_points:
                            translated_point = self.translate_text(point, target_lang=target_language, source_lang='hi')
                            translated_points.append(translated_point)
                        key_points = translated_points
                    except Exception as e:
                        print(f"Error translating Hindi summary: {e}")

            else:
                # For non-Hindi documents
                if target_language == 'hi':
                    # For English docs with Hindi target, first generate in English then translate
                    summary = self._generate_llm_summary(content, language='en')
                    key_points = self._extract_llm_key_points(content, language='en')

                    # Translate to Hindi
                    try:
                        summary = self.translate_text(summary, target_lang='hi', source_lang='en')
                        translated_points = []
                        for point in key_points:
                            translated_point = self.translate_text(point, target_lang='hi', source_lang='en')
                            translated_points.append(translated_point)
                        key_points = translated_points
                    except Exception as e:
                        print(f"Error translating summary to Hindi: {e}")
                else:
                    # Both source and target are English
                    summary = self._generate_llm_summary(content, language='en')
                    key_points = self._extract_llm_key_points(content, language='en')
        else:
            # Use extractive summarization as fallback
            print("Using extractive summarization")
            summary = self._generate_text_summary(content)
            key_points = self._extract_key_points(content)

            # If languages don't match, translate the extractive summary
            if doc_lang != target_language:
                try:
                    summary = self.translate_text(summary, target_lang=target_language, source_lang=doc_lang)
                    translated_points = []
                    for point in key_points:
                        translated_point = self.translate_text(point, target_lang=target_language, source_lang=doc_lang)
                        translated_points.append(translated_point)
                    key_points = translated_points
                except Exception as e:
                    print(f"Error translating extractive summary: {e}")

        return {
            "summary": summary,
            "key_points": key_points,
            "language": target_language
        }

    def get_related_documents(self, doc_id: str, k: int = 5) -> List[Dict[str, Any]]:
        """
        Find documents that are similar/related to the given document
        Returns a list of document summaries ordered by relevance
        """
        # Get the source document
        source_chunks = self.get_document(doc_id)
        if not source_chunks:
            return []

        # Get combined text from first few chunks of source document
        source_text = "\n\n".join([chunk.page_content for chunk in source_chunks[:3]])
        source_lang = source_chunks[0].metadata.get('language', 'en')

        # Use source document content as a query to find similar documents
        results = self.search_documents(source_text, k=k + 1, query_language=source_lang)

        # Filter out the source document itself
        results = [r for r in results if r['doc_id'] != doc_id]

        # Return the top k results
        return results[:k]

    def find_related_by_topic(self, topic: str, k: int = 5, exclude_doc_id: str = None) -> List[Dict[str, Any]]:
        """
        Find documents related to a specific topic
        Useful for exploring document collections by subject
        """
        # Detect the language of the topic
        topic_lang = self.detect_language(topic)

        # Search for documents related to this topic
        results = self.search_documents(topic, k=k + 1 if exclude_doc_id else k, query_language=topic_lang)

        # Filter out the excluded document if specified
        if exclude_doc_id:
            results = [r for r in results if r['doc_id'] != exclude_doc_id]

        # Return the top k results
        return results[:k]

    def get_document_recommendations(self, doc_id: str, user_query: str = "", k: int = 3) -> List[Dict[str, Any]]:
        """
        Get document recommendations based on both document content and user query
        This combines content-based and query-based recommendations
        """
        results = []

        # Get document-based recommendations
        doc_based_results = self.get_related_documents(doc_id, k=max(2, k))

        # If we have a user query, also get query-based recommendations
        if user_query:
            query_lang = self.detect_language(user_query)
            query_based_results = self.search_documents(user_query, k=max(2, k), query_language=query_lang)
            query_based_results = [r for r in query_based_results if r['doc_id'] != doc_id]
        else:
            query_based_results = []

        # If we have both types of results, combine them with priority to query-based ones
        if query_based_results and doc_based_results:
            # Start with query-based results
            results.extend(query_based_results)

            # Add document-based results not already included
            doc_ids_added = set(r['doc_id'] for r in results)
            for result in doc_based_results:
                if result['doc_id'] not in doc_ids_added:
                    results.append(result)
                    doc_ids_added.add(result['doc_id'])
                    if len(results) >= k:
                        break
        # If we only have one type, use those
        elif query_based_results:
            results = query_based_results
        else:
            results = doc_based_results

        # Limit to the requested number of results
        return results[:k]

    def extract_top_topics(self, k: int = 5) -> List[Dict[str, Any]]:
        """
        Extract the top k topics from the entire document collection
        Useful for creating exploration points or document categorization
        """
        if not self.documents:
            return []

        # Combine sample text from all documents
        all_text = ""
        for doc_id, chunks in self.documents.items():
            if chunks:
                # Take first chunk from each document
                sample = chunks[0].page_content[:500]  # First 500 chars
                all_text += sample + "\n\n"

        # Simple topic extraction
        topics = []

        # Use regex to extract important noun phrases
        import re
        # Find capitalized phrases that might represent topics
        topic_matches = re.findall(r'\b([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)\b', all_text)
        topic_counts = {}

        for match in topic_matches:
            if len(match.split()) >= 1:  # At least one word
                topic_counts[match] = topic_counts.get(match, 0) + 1

        # Sort by frequency
        sorted_topics = sorted(topic_counts.items(), key=lambda x: x[1], reverse=True)

        # Create the topic results
        for topic, count in sorted_topics[:k]:
            # Find a sample document for this topic
            sample_docs = self.find_related_by_topic(topic, k=1)

            topics.append({
                'topic': topic,
                'count': count,
                'sample_doc': sample_docs[0]['doc_id'] if sample_docs else None
            })

        return topics

    def fetch_document_from_url(self, url: str) -> str:
        """
        Fetch documents from online sources using robust error handling
        Returns a document ID if successful
        """
        try:
            # Create a temp directory if needed
            temp_dir = os.path.join(self.storage_dir, "temp")
            os.makedirs(temp_dir, exist_ok=True)

            # Extract filename from URL or create a generic one
            filename = os.path.basename(url)
            if not filename or '.' not in filename:
                # Determine extension based on URL pattern
                if 'pdf' in url.lower():
                    filename = f"document_{int(time.time())}.pdf"
                elif 'doc' in url.lower():
                    filename = f"document_{int(time.time())}.docx"
                else:
                    filename = f"document_{int(time.time())}.txt"

            # Download the file with a timeout
            response = requests.get(url, timeout=30, stream=True)
            response.raise_for_status()  # Raise exception for HTTP errors

            # Save the file locally
            file_path = os.path.join(temp_dir, filename)
            with open(file_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            # Process the downloaded document
            doc_id = self.process_document(file_path, filename)
            return doc_id

        except requests.exceptions.RequestException as e:
            print(f"Error fetching document from URL: {e}")
            raise ValueError(f"Could not download document: {e}")
        except Exception as e:
            print(f"Error processing online document: {e}")
            raise ValueError(f"Error processing document: {e}")